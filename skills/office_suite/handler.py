import os
import structlog
from typing import Dict, Any, List
# Import libraries inside methods or try/except to avoid crashing if not installed yet
try:
    from docx import Document
    from docx.shared import Pt
    from pptx import Presentation
    from pptx.util import Inches, Pt as PptPt
except ImportError:
    pass # Will be installed in Docker

from skills.base import BaseSkill, SkillMetadata
from skills.office_suite.schema import GenerateDocument, DocumentElement

logger = structlog.get_logger(__name__)

class OfficeSkill(BaseSkill):
    name = "office_suite"
    description = "Generates professional Office documents (Word .docx, PowerPoint .pptx)."

    @property
    def metadata(self) -> SkillMetadata:
        return SkillMetadata(
            capabilities=["document_generation", "office", "word", "powerpoint"],
            risk_level="low",
            side_effects=True,
            idempotent=True,
            requires_network=False,
            requires_filesystem=True,
            cost_class="medium"
        )

    def get_system_prompt(self) -> str:
        return """
        You are an expert Office Assistant capable of generating professional documents and presentations.
        
        When a user asks for a "report", "document", "TZ", or "specs" -> Use 'docx' mode.
        When a user asks for a "presentation", "slides", or "deck" -> Use 'pptx' mode.
        
        Structure your content logically:
        - Use 'heading' for section titles.
        - Use 'paragraph' for body text.
        - Use 'bullet_list' for key points.
        
        For Presentations:
        - Keep content concise.
        - Each 'heading' usually starts a new slide (implied logic).
        """

    @property
    def input_schema(self) -> Dict[str, Any]:
        return GenerateDocument

    def is_sensitive(self, params: Any) -> bool:
        """Creating files is a sensitive operation."""
        return True

    async def execute(self, params: Any, state: Any = None) -> str:
        try:
            # Handle both Dict and Pydantic object
            if isinstance(params, dict):
                req = GenerateDocument(**params)
            else:
                req = params
            
            output_dir = "generated_files"
            os.makedirs(output_dir, exist_ok=True)
            
            filename = f"{req.file_name}.{req.doc_type}"
            filepath = os.path.join(output_dir, filename)

            if req.doc_type == "docx":
                self._generate_docx(req, filepath)
            elif req.doc_type == "pptx":
                self._generate_pptx(req, filepath)
            else:
                return f"Error: Unsupported doc_type: {req.doc_type}"
                
            return f"✅ Document generated successfully: [{filename}](file:///{filepath.replace(os.sep, '/')})"

        except Exception as e:
            logger.error("office_skill_error", error=str(e))
            return {"error": f"Failed to generate document: {str(e)}"}

    def _generate_docx(self, req: GenerateDocument, filepath: str):
        from docx import Document
        doc = Document()
        
        # Title
        doc.add_heading(req.title, 0)
        
        for el in req.elements:
            if el.type == "heading":
                doc.add_heading(el.content, level=1)
            elif el.type == "title":
                 doc.add_heading(el.content, level=0)
            elif el.type == "bullet_list":
                doc.add_paragraph(el.content, style='List Bullet')
            elif el.type == "numbered_list":
                doc.add_paragraph(el.content, style='List Number')
            elif el.type == "code_block":
                p = doc.add_paragraph(el.content)
                p.style = 'Quote' # Fallback style for code
            elif el.type == "image":
                 # Check if file exists, or check in generated_files
                 img_path = el.content
                 if not os.path.exists(img_path):
                     candidate = os.path.join("generated_files", img_path)
                     if os.path.exists(candidate):
                         img_path = candidate
                 
                 if os.path.exists(img_path):
                     try:
                         doc.add_picture(img_path, width=Inches(6.0))
                     except Exception as e:
                         doc.add_paragraph(f"[Image Error: {el.content} - {str(e)}]")
                 else:
                     doc.add_paragraph(f"[Image not found: {el.content}]")
            else:
                # Default paragraph
                doc.add_paragraph(el.content)
                
        doc.save(filepath)

    def _generate_pptx(self, req: GenerateDocument, filepath: str):
        from pptx import Presentation
        prs = Presentation()
        
        # Title Slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = req.title
        subtitle.text = "Generated by SGR Core Agent"
        
        # Content Slides
        bullet_slide_layout = prs.slide_layouts[1]
        
        current_slide = None
        current_body = None
        
        for el in req.elements:
            # Logic: Heading starts a new slide
            if el.type in ["heading", "title"] or current_slide is None:
                current_slide = prs.slides.add_slide(bullet_slide_layout)
                shapes = current_slide.shapes
                title_shape = shapes.title
                title_shape.text = el.content if el.type in ["heading", "title"] else "Slide"
                current_body = shapes.placeholders[1].text_frame
                
                # If the element was just a heading, we are done for this iteration
                if el.type in ["heading", "title"]:
                    continue

            # Add content to current slide
            if el.type == "bullet_list":
                p = current_body.add_paragraph()
                p.text = el.content
                p.level = 0
            elif el.type == "paragraph":
                p = current_body.add_paragraph()
                p.text = el.content
                p.level = 0
            elif el.type == "image":
                # Check if file exists, or check in generated_files
                img_path = el.content
                if not os.path.exists(img_path):
                     candidate = os.path.join("generated_files", img_path)
                     if os.path.exists(candidate):
                         img_path = candidate

                if os.path.exists(img_path):
                    try:
                        # Add image to current slide
                        # Position: right side or centered safely
                        # Using Inces from pptx.util
                        top = Inches(2)
                        left = Inches(1)
                        height = Inches(4)
                        current_slide.shapes.add_picture(img_path, left, top, height=height)
                    except Exception as e:
                         # Fallback if image fails
                         p = current_body.add_paragraph()
                         p.text = f"[Image Error: {el.content}]"
                else:
                     p = current_body.add_paragraph()
                     p.text = f"[Image not found: {el.content}]"

        prs.save(filepath)
